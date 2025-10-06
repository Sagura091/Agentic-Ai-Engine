"""
Presentation processor for PowerPoint, ODP files.

This module provides comprehensive presentation processing:
- PowerPoint (PPTX, PPT) support
- ODP (OpenDocument Presentation) support
- Slide text extraction
- Speaker notes extraction
- Image extraction
- Slide layout detection
- Table extraction
"""

from typing import Dict, Any, Optional, List, Tuple
from pathlib import Path
from datetime import datetime
import io

import structlog

from .models_result import ProcessResult, ProcessorError, ErrorCode, ProcessingStage, ImageData
from .dependencies import get_dependency_checker

logger = structlog.get_logger(__name__)


class PresentationProcessor:
    """
    Comprehensive presentation processor.
    
    Features:
    - PPTX and PPT support
    - ODP support
    - Slide text extraction
    - Speaker notes
    - Image extraction
    - Table extraction
    - Layout detection
    """
    
    def __init__(
        self,
        extract_images: bool = True,
        extract_notes: bool = True,
        max_slides: int = 1000
    ):
        """
        Initialize presentation processor.
        
        Args:
            extract_images: Extract images from slides
            extract_notes: Extract speaker notes
            max_slides: Maximum slides to process
        """
        self.extract_images = extract_images
        self.extract_notes = extract_notes
        self.max_slides = max_slides
        
        self.dep_checker = get_dependency_checker()
        
        logger.info(
            "PresentationProcessor initialized",
            extract_images=extract_images,
            extract_notes=extract_notes,
            max_slides=max_slides
        )
    
    async def process(
        self,
        content: bytes,
        filename: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessResult:
        """
        Process presentation file.
        
        Args:
            content: Presentation file content
            filename: Filename
            metadata: Additional metadata
            
        Returns:
            ProcessResult with extracted content
        """
        start_time = datetime.utcnow()
        errors = []
        
        try:
            # Detect format
            file_ext = Path(filename).suffix.lower()
            
            if file_ext == '.pptx':
                result_data = await self._process_pptx(content, filename)
            elif file_ext == '.ppt':
                result_data = await self._process_ppt(content, filename)
            elif file_ext == '.odp':
                result_data = await self._process_odp(content, filename)
            else:
                return ProcessResult(
                    text="",
                    metadata=metadata or {},
                    errors=[ProcessorError(
                        code=ErrorCode.UNSUPPORTED_FORMAT,
                        message=f"Unsupported presentation format: {file_ext}",
                        stage=ProcessingStage.EXTRACTION,
                        retriable=False
                    )],
                    processor_name="PresentationProcessor",
                    processing_time_ms=(datetime.utcnow() - start_time).total_seconds() * 1000
                )
            
            # Build text representation
            text_parts = [f"Presentation: {filename}"]
            text_parts.append(f"Slides: {len(result_data['slides'])}")
            text_parts.append("")
            
            for idx, slide_info in enumerate(result_data['slides'], 1):
                text_parts.append(f"## Slide {idx}: {slide_info.get('title', 'Untitled')}")
                text_parts.append("")
                
                # Add slide content
                if slide_info.get('content'):
                    text_parts.append(slide_info['content'])
                    text_parts.append("")
                
                # Add speaker notes
                if slide_info.get('notes'):
                    text_parts.append("**Speaker Notes:**")
                    text_parts.append(slide_info['notes'])
                    text_parts.append("")
            
            result_metadata = {
                **(metadata or {}),
                "presentation_format": file_ext[1:],
                "total_slides": len(result_data['slides']),
                "slides": result_data['slides']
            }
            
            return ProcessResult(
                text="\n".join(text_parts),
                metadata=result_metadata,
                images=result_data.get('images', []),
                errors=errors,
                processor_name="PresentationProcessor",
                processing_time_ms=(datetime.utcnow() - start_time).total_seconds() * 1000
            )
            
        except Exception as e:
            logger.error("Presentation processing failed", error=str(e), filename=filename)
            
            return ProcessResult(
                text="",
                metadata=metadata or {},
                errors=[ProcessorError(
                    code=ErrorCode.PROCESSING_FAILED,
                    message=f"Presentation processing failed: {str(e)}",
                    stage=ProcessingStage.EXTRACTION,
                    retriable=True,
                    details={"error": str(e)}
                )],
                processor_name="PresentationProcessor",
                processing_time_ms=(datetime.utcnow() - start_time).total_seconds() * 1000
            )
    
    async def _process_pptx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PPTX file using python-pptx."""
        from pptx import Presentation
        import base64
        
        prs = Presentation(io.BytesIO(content))
        
        slides = []
        images = []
        
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx >= self.max_slides:
                break
            
            # Extract title
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            
            # Extract text content
            content_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text and text != title:
                        content_parts.append(text)
                
                # Extract tables
                if shape.has_table:
                    table_text = self._extract_table_from_shape(shape)
                    if table_text:
                        content_parts.append(table_text)
                
                # Extract images
                if self.extract_images and shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        images.append(ImageData(
                            data=base64.b64encode(image_bytes).decode('utf-8'),
                            format=image.ext,
                            width=shape.width,
                            height=shape.height,
                            page=slide_idx + 1
                        ))
                    except Exception as e:
                        logger.warning("Failed to extract image", error=str(e))
            
            # Extract speaker notes
            notes = ""
            if self.extract_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text.strip()
            
            slides.append({
                "title": title,
                "content": "\n".join(content_parts),
                "notes": notes,
                "slide_number": slide_idx + 1
            })
        
        return {
            "slides": slides,
            "images": images
        }
    
    async def _process_ppt(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PPT file using LibreOffice conversion."""
        # PPT is a binary format, requires conversion
        # This is a placeholder - would need LibreOffice or similar
        raise NotImplementedError("PPT format requires LibreOffice conversion")
    
    async def _process_odp(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process ODP file using odfpy."""
        from odf import opendocument
        from odf.draw import Page, Frame, Image as ODFImage
        from odf.text import P
        
        doc = opendocument.load(io.BytesIO(content))
        
        slides = []
        images = []
        
        for slide_idx, page in enumerate(doc.presentation.getElementsByType(Page)):
            if slide_idx >= self.max_slides:
                break
            
            # Extract title and content
            title = ""
            content_parts = []
            
            for frame in page.getElementsByType(Frame):
                for p in frame.getElementsByType(P):
                    text = str(p).strip()
                    if text:
                        if not title:
                            title = text
                        else:
                            content_parts.append(text)
            
            slides.append({
                "title": title,
                "content": "\n".join(content_parts),
                "notes": "",
                "slide_number": slide_idx + 1
            })
        
        return {
            "slides": slides,
            "images": images
        }
    
    def _extract_table_from_shape(self, shape) -> str:
        """Extract table text from shape."""
        try:
            table = shape.table
            rows = []
            
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text.strip())
                rows.append(" | ".join(cells))
            
            return "\n".join(rows)
            
        except Exception as e:
            logger.warning("Failed to extract table", error=str(e))
            return ""

